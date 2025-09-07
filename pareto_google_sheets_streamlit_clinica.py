# app.py
import io
import re
import zipfile
from collections import Counter

import numpy as np
import pandas as pd
import streamlit as st
from plotly.subplots import make_subplots
import plotly.graph_objects as go
import importlib.util

# =====================
# CONFIGURAÇÃO INICIAL
# =====================
st.set_page_config(layout="wide", page_title="Gráficos de Pareto")
st.title("📊 Gráficos de Pareto")
st.caption("Fonte: Respostas da planilha do Google Sheets (atualiza em tempo real)")

# CSS básico
st.markdown("""
<style>
.block-container { padding-top: 1rem; padding-bottom: 2rem; }
[data-testid="column"] { padding: 0.25rem; }
.stPlotlyChart, .stImage, .stMarkdown img { width: 100% !important; }
</style>
""", unsafe_allow_html=True)

# Checagem do Kaleido (para exportar imagens)
KALEIDO_OK = importlib.util.find_spec("kaleido") is not None

# =====================
# FUNÇÕES AUXILIARES
# =====================
def slugify(texto: str) -> str:
    s = re.sub(r"\s+", "_", texto.strip())
    s = re.sub(r"[^\w\-_.()]+", "", s, flags=re.UNICODE)
    return s.lower()

def carregar_planilha_google_sheets(url: str):
    try:
        if "docs.google.com/spreadsheets/d/e/" in url:
            published_id = url.split("/d/e/")[1].split("/")[0]
            csv_url = f"https://docs.google.com/spreadsheets/d/e/{published_id}/pub?output=csv"
        elif "docs.google.com/spreadsheets/d/" in url:
            sheet_id = url.split("/d/")[1].split("/")[0]
            csv_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv"
        else:
            st.error("URL inválida. Verifique se é um link público do Google Sheets.")
            return None
        return pd.read_csv(csv_url)
    except Exception as e:
        st.error(f"Erro ao carregar dados: {e}")
        return None

def contar_respostas_multipla(df: pd.DataFrame, coluna: str) -> Counter:
    todas_respostas = []
    for resposta in df[coluna].dropna():
        opcoes = [op.strip() for op in str(resposta).split(",")]
        todas_respostas.extend(opcoes)
    return Counter(todas_respostas)

def criar_figura_pareto_plotly(counter: Counter, titulo: str):
    if not counter:
        fig = make_subplots(specs=[[{"secondary_y": True}]])
        fig.add_annotation(text="Sem dados para exibir", x=0.5, y=0.5, showarrow=False)
        fig.update_layout(margin=dict(l=20,r=20,t=80,b=20), height=360)
        return fig

    labels, valores = zip(*counter.most_common())
    totais = np.array(valores, dtype=float)
    acumulado = np.cumsum(totais)
    p_acum = 100 * acumulado / acumulado[-1]

    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(go.Bar(x=list(labels), y=list(totais), name="Frequência"), secondary_y=False)
    fig.add_trace(go.Scatter(x=list(labels), y=list(p_acum), mode="lines+markers", name="% Acumulado"), secondary_y=True)
    fig.add_hline(y=80, line_dash="dash", line_color="gray", secondary_y=True)

    fig.update_yaxes(title_text="Frequência", secondary_y=False)
    fig.update_yaxes(title_text="% Acumulado", range=[0,110], secondary_y=True)
    fig.update_xaxes(tickangle=45)

    # TÍTULO “no alto”, como header
    fig.update_layout(
        title={"text": titulo, "x": 0.5, "y": 0.98, "xanchor": "center", "yanchor": "top"},
        margin=dict(l=20, r=20, t=90, b=40),
        height=380,
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
    )
    return fig

def fig_to_bytes(fig, fmt="png", scale=2):
    if not KALEIDO_OK:
        return None
    return fig.to_image(format=fmt, scale=scale, engine="kaleido")

def build_zip(images: list[tuple[str, bytes]]) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for nome, data in images:
            zf.writestr(nome, data)
    buf.seek(0)
    return buf.read()

def build_pptx(pares: list[tuple[str, bytes]]):
    """Gera PPTX com 1 slide por gráfico (título + imagem PNG)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # Layout título + conteúdo (geralmente index 1)
    layout = prs.slide_layouts[1]

    for titulo, png_bytes in pares:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = titulo

        # Caixa de conteúdo (placeholder 1)
        ph = slide.placeholders[1]
        # Inserir figura dimensionando para largura do placeholder
        pic = ph.insert_picture(io.BytesIO(png_bytes))

        # Opcional: centralizar título
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

# =====================
# URL DA PLANILHA
# =====================
URL_GOOGLE_SHEETS = "https://docs.google.com/spreadsheets/d/e/2PACX-1vSKwCflZovzD_0UAlLsDplKqWz2-WKs3agK-HaDQFmT6jx9RkkUAXNbUJvsD622uqUWTpXUTQ8XgILV/pub?output=csv"

# =====================
# CARREGAR DADOS
# =====================
df = carregar_planilha_google_sheets(URL_GOOGLE_SHEETS)

if df is not None:
    st.success("✅ Planilha carregada com sucesso.")
    st.caption(f"Linhas carregadas: {len(df)}")

    perguntas = df.columns[1:8]

    # Config da barra do Plotly (ícone câmera → 1920x1080)
    plotly_config = {
        "displayModeBar": True,
        "scrollZoom": True,
        "responsive": True,
        "toImageButtonOptions": {
            "format": "png",
            "filename": "grafico",
            "height": 1080,
            "width": 1920,
            "scale": 1
        }
    }

    # Guardar figs e imagens para ZIP/PPTX
    figs: list[tuple[str, go.Figure]] = []
    pngs: list[tuple[str, bytes]] = []

    for i in range(0, len(perguntas), 2):
        colunas = st.columns(2)
        for j in range(2):
            if i + j < len(perguntas):
                coluna = perguntas[i + j]
                titulo = f"{i + j + 1}) {coluna}"
                contador = contar_respostas_multipla(df, coluna)
                fig = criar_figura_pareto_plotly(contador, titulo)

                with colunas[j]:
                    st.plotly_chart(fig, use_container_width=True, config=plotly_config)

                    fname_base = slugify(titulo)

                    # Botões individuais
                    c1, c2 = st.columns([1,1])
                    with c1:
                        if KALEIDO_OK:
                            png_bytes = fig_to_bytes(fig, fmt="png", scale=2)
                            st.download_button(
                                "Baixar PNG (alta)",
                                data=png_bytes,
                                file_name=f"{fname_base}.png",
                                mime="image/png",
                                key=f"png_{i}_{j}",
                            )
                        else:
                            st.info("Instale 'kaleido' para habilitar o download PNG/SVG.")
                    with c2:
                        if KALEIDO_OK:
                            svg_bytes = fig_to_bytes(fig, fmt="svg", scale=1)
                            st.download_button(
                                "Baixar SVG (vetor)",
                                data=svg_bytes,
                                file_name=f"{fname_base}.svg",
                                mime="image/svg+xml",
                                key=f"svg_{i}_{j}",
                            )

                # Acumula para ZIP/PPTX
                figs.append((titulo, fig))
                if KALEIDO_OK:
                    png_bytes = fig_to_bytes(fig, fmt="png", scale=2)
                    if png_bytes:
                        pngs.append((f"{slugify(titulo)}.png", png_bytes))

    st.markdown("---")

    # Ações globais (rodapé)
    colA, colB = st.columns([1,1])

    with colA:
        if KALEIDO_OK and pngs:
            zip_bytes = build_zip(pngs)
            st.download_button(
                "⬇️ Baixar TODOS (ZIP, PNG alta)",
                data=zip_bytes,
                file_name="graficos_pareto.zip",
                mime="application/zip",
                key="zip_all",
            )
        else:
            st.caption("Para baixar todos: adicione 'kaleido' no requirements.txt.")

    with colB:
        if KALEIDO_OK and pngs:
            try:
                pptx_bytes = build_pptx([(name[:-4].replace("_", " "), data) for name, data in pngs])
                st.download_button(
                    "🖼️ Gerar PPTX (1 slide por gráfico)",
                    data=pptx_bytes,
                    file_name="graficos_pareto.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="pptx_all",
                )
            except Exception as e:
                st.warning(f"Não foi possível gerar o PPTX: {e}")
        else:
            st.caption("Para gerar PPTX: adicione 'kaleido' e 'python-pptx' ao requirements.")

else:
    st.error("❌ Não foi possível carregar os dados da planilha.")
